# backend/pptx_generation.py
# 使用 python-pptx 库生成 .pptx 文件，优先插入本地图片

import os
import requests
import tempfile
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.shapes.placeholder import PP_PLACEHOLDER # <--- 确保这里导入的是 PP_PLACEHOLDER_TYPE

# --- 简单的样式配置 (MVP阶段，后续通过模板管理) ---
TITLE_FONT_SIZE = Pt(30)
BODY_FONT_SIZE = Pt(18)
NOTES_FONT_SIZE = Pt(10)
BULLET_CHAR = '•' # 或 '-'

# --- 辅助函数：下载图片到临时文件 (保留，作为未来外部图片的可能性) ---
# 现在优先使用本地图片，这个函数作为 fallback 或 future work
def download_image_to_temp(image_url: str) -> str | None:
    """
    从给定的 URL 下载图片并保存到临时文件。
    返回临时文件的路径，如果下载失败则返回 None。
    """
    # ... (函数内容与之前相同，无需修改) ...
    try:
        print(f"Attempting to download image from: {image_url}")
        response = requests.get(image_url, stream=True)
        response.raise_for_status()

        content_type = response.headers.get('Content-Type', '').lower()
        if 'jpeg' in content_type or 'jpg' in content_type:
            ext = '.jpg'
        elif 'png' in content_type:
            ext = '.png'
        elif 'gif' in content_type:
            ext = '.gif'
        else:
            print(f"Warning: Unknown image content type: {content_type}. Assuming .jpg")
            ext = '.jpg'

        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
            for chunk in response.iter_content(chunk_size=8192):
                tmp_file.write(chunk)
            temp_path = tmp_file.name

        print(f"Image downloaded successfully to temporary file: {temp_path}")
        return temp_path

    except requests.exceptions.RequestException as e:
        print(f"Error downloading image from {image_url}: {e}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during image download: {e}")
        return None


# --- PPTX 生成核心函数 ---
# 接收包含 local_image_path 的 slides_data
def create_presentation_file(presentation_id: str, slides_data: list, template_path: str | None = None):
    """
    根据幻灯片数据生成 PPTX 文件，可选择加载模板，并尝试插入本地图片。

    Args:
        presentation_id (str): 演示文稿的唯一 ID。
        slides_data (list): 包含每个幻灯片数据的列表（每个dict可能包含 local_image_path）。
        template_path (str | None): 可选的模板文件物理路径。

    Returns:
        str: 生成的 PPTX 文件的完整文件路径。
    """
    if template_path and os.path.exists(template_path):
        try: # 尝试加载模板，如果模板文件无效会在这里抛出异常
             prs = PPTXPresentation(template_path)
             print(f"Loaded presentation template from: {template_path}")
        except Exception as e:
             print(f"Error loading template {template_path}: {e}")
             print("Falling back to blank presentation.")
             prs = PPTXPresentation() # 加载失败则创建空白演示文稿
    else:
        prs = PPTXPresentation()
        print("Created a blank presentation (no valid template found or specified).")


    # --- 幻灯片布局选择 ---
    title_content_layout = None
    blank_layout = None

    for layout in prs.slide_layouts:
        has_content_placeholder = any(p.is_placeholder and p.placeholder_format.idx == 1 for p in layout.placeholders)
        is_blank_layout = len(layout.placeholders) == 0

        if has_content_placeholder:
             title_content_layout = layout
             break

        if is_blank_layout:
             blank_layout = layout


    if not title_content_layout:
        if blank_layout:
             title_content_layout = blank_layout
             print("Warning: 'Title and Content' layout not found, using a blank layout.")
        elif prs.slide_layouts:
             title_content_layout = prs.slide_layouts[0]
             print("Warning: No suitable layout found, using the first available layout.")
        else:
             raise RuntimeError("No usable slide layout found in the presentation or template.")


    # --- 遍历幻灯片数据，创建每一页幻灯片 ---
    for slide_info in slides_data:
        slide = prs.slides.add_slide(title_content_layout)

        # --- 添加标题 ---
        title_placeholder = None
        for shape in slide.placeholders:
             # 使用 PP_PLACEHOLDER_TYPE 查找标题占位符
             if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                  title_placeholder = shape
                  break
             if shape.is_placeholder and shape.placeholder_format.idx == 0:
                  title_placeholder = shape
                  break


        if title_placeholder:
             title_placeholder.text = slide_info.get("title", "Untitled Slide")
        else:
             left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(1)
             title_textbox = slide.shapes.add_textbox(left, top, width, height)
             tf = title_textbox.text_frame
             p = tf.add_paragraph()
             p.text = slide_info.get("title", "Untitled Slide")
             font = p.runs[0].font
             font.size = TITLE_FONT_SIZE
             font.bold = True
             print(f"Warning: No title placeholder found for slide. Added title as a textbox.")


        # --- 添加内容要点 (文本框) ---
        content_placeholder = None
        for shape in slide.placeholders:
             # 使用 PP_PLACEHOLDER_TYPE 查找内容或对象占位符
             if shape.is_placeholder and shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                  content_placeholder = shape
                  break
             if shape.is_placeholder and shape.placeholder_format.idx == 1:
                  content_placeholder = shape
                  break


        if content_placeholder:
             tf = content_placeholder.text_frame
             tf.clear()

             points = slide_info.get("points", [])

             if points:
                 p = tf.add_paragraph()
                 p.text = points[0]
                 p.level = 0

                 for point in points[1:]:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0

        else:
             left, top, width, height = Inches(0.5), Inches(1.5), Inches(6), Inches(5)
             content_textbox = slide.shapes.add_textbox(left, top, width, height)
             tf = content_textbox.text_frame
             points = slide_info.get("points", [])
             if points:
                 for point in points:
                      p = tf.add_paragraph()
                      p.text = point
                      p.level = 0
             print(f"Warning: No content placeholder found for slide. Added content as a textbox.")


        # --- 添加图片 ---
        # 优先使用匹配到的本地图片路径
        local_image_relative_path = slide_info.get("local_image_path")
        image_inserted = False # 标记是否已成功插入图片

        if local_image_relative_path:
            # 构建本地图片的完整物理路径
            # 需要知道静态文件目录的绝对路径，可以在 main.py 中传递进来，或者在这里重新构建
            # 为了简化，我们在 main.py 中构建完整路径并传递进来更方便
            # TODO: main.py 需要修改，将 local_image_path 转为完整物理路径传递给 create_presentation_file
            # 暂时先假设传递进来的是完整路径 for testing this file in isolation
            # 或者在 main.py 中处理，这里只接收完整路径

            # *** 重要修正：create_presentation_file 需要接收图片的完整物理路径 ***
            # 让我们修改函数签名，接收 local_image_full_path
            # def create_presentation_file(..., local_image_full_path: str | None = None):
            # 然后在 main.py 中调用时，根据 local_image_relative_path 构建完整路径并传递
            # 现在先按接收相对路径并在后端构建完整路径的方式来写，但这需要知道 STATIC_DIR
            # 更干净是在 main.py 中处理路径
            # For now, let's assume local_image_relative_path is actually the full path passed from main.py
            local_image_full_path = local_image_relative_path # <--- 临时假设 local_image_relative_path 是完整路径

            if local_image_full_path and os.path.exists(local_image_full_path):
                try:
                    # 尝试找到图片占位符并插入
                    image_placeholder = None
                    for shape in slide.placeholders:
                        # 使用 PP_PLACEHOLDER_TYPE 查找图片占位符
                        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                            image_placeholder = shape
                            break

                    if image_placeholder:
                        # 如果找到了图片占位符
                        image_placeholder.insert_picture(local_image_full_path)
                        print(f"Inserted local image '{local_image_full_path}' into picture placeholder on slide.")
                        image_inserted = True
                    else:
                        # 如果没有找到图片占位符，添加到右侧中间固定位置
                        left = Inches(7)
                        top = Inches(2)
                        width = Inches(2.5)
                        slide.shapes.add_picture(local_image_full_path, left, top, width=width)
                        print(f"No picture placeholder found. Added local image '{local_image_full_path}' to a fixed position on slide.")
                        image_inserted = True

                except Exception as e:
                     print(f"Error inserting local image {local_image_full_path} into slide: {e}")
                # No need to delete temp file here, it's a local static file

        # --- 如果本地图片未插入，尝试使用外部 URL (保留作为 Fallback 或未来可能) ---
        # 在当前 MVP 中，ai_generation 不再生成 image_url，所以这个分支不会被触发
        # 但保留代码结构，以便未来可以轻松切换回外部图片源
        if not image_inserted:
             image_url = slide_info.get("image_url") # 这个 image_url 现在应该总是 None
             if image_url:
                temp_image_path = download_image_to_temp(image_url)

                if temp_image_path and os.path.exists(temp_image_path):
                    try:
                         image_placeholder = None # 重新查找占位符，避免冲突
                         for shape in slide.placeholders:
                             if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                                 image_placeholder = shape
                                 break

                         if image_placeholder:
                             image_placeholder.insert_picture(temp_image_path)
                             print(f"Inserted downloaded image from {image_url} into picture placeholder on slide.")
                             image_inserted = True # 标记已插入
                         else:
                             left = Inches(7)
                             top = Inches(2)
                             width = Inches(2.5)
                             slide.shapes.add_picture(temp_image_path, left, top, width=width)
                             print(f"No picture placeholder found. Added downloaded image from {image_url} to a fixed position on slide.")
                             image_inserted = True # 标记已插入

                    except Exception as e:
                         print(f"Error inserting downloaded image {temp_image_path} into slide: {e}")
                    finally:
                         # 清理临时文件
                         try:
                             os.remove(temp_image_path)
                             print(f"Removed temporary image file: {temp_image_path}")
                         except OSError as e:
                             print(f"Error removing temporary image file {temp_image_path}: {e}")
                else:
                     print(f"Warning: Could not download or find temporary image file for URL: {image_url}. Image not inserted.")


        # --- 添加演讲者备注 ---
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.clear()
        text_frame.text = slide_info.get("notes", "")


    # --- 保存 PPTX 文件 ---
    output_dir = "generated_pptx"
    os.makedirs(output_dir, exist_ok=True)
    filepath = os.path.join(output_dir, f"presentation_{presentation_id}.pptx")
    prs.save(filepath)

    return filepath

# --- 用于本地测试此模块的代码 (可选) ---
if __name__ == "__main__":
    # 假设的幻灯片数据 (包含 local_image_path)
    # 需要先在 backend/static/images 目录下放入 tech1.jpg, city1.jpg 等文件
    import sys
    # 将 backend 目录添加到 sys.path，以便导入 ai_generation 中的 LOCAL_IMAGES_DIR
    sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    from ai_generation import LOCAL_IMAGES_DIR # 导入 LOCAL_IMAGES_DIR

    # 在 backend/static/images 目录下创建一个虚拟的 tech1.jpg 文件用于测试
    test_image_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'static', 'images')
    os.makedirs(test_image_dir, exist_ok=True)
    # 创建一个空白图片文件
    try:
         from PIL import Image # 需要安装 Pillow 库: pip install Pillow
         img = Image.new('RGB', (60, 30), color = (255, 0, 0))
         test_image_path = os.path.join(test_image_dir, 'tech1.jpg')
         img.save(test_image_path)
         print(f"Created test image file: {test_image_path}")
         test_local_image_path = test_image_path # 直接使用完整路径进行测试
         test_image_present = True
    except ImportError:
         print("Pillow not installed. Cannot create test image. Image insertion test will likely fail.")
         test_local_image_path = None
         test_image_present = False
    except Exception as e:
         print(f"Error creating test image: {e}")
         test_local_image_path = None
         test_image_present = False


    test_slides_data_with_local_image = [
        {
            "title": "测试本地图片插入",
            "points": ["这是一个要点", "本地图片"],
            "notes": "这页会尝试插入本地图片。",
            # 传递本地图片的完整物理路径给 create_presentation_file
            "local_image_path": test_local_image_path if test_image_present else None
        },
         {
            "title": "无图片页",
            "points": ["没有图片的要点"],
            "notes": "这页没有图片。",
            "local_image_path": None
        }
    ]
    test_id = "test_local_image_pptx"

    # 需要指定一个实际存在的模板文件路径进行测试，或者设置为 None 使用空白演示文稿
    BASE_DIR_PPPTX = os.path.dirname(os.path.abspath(__file__)) # 当前文件目录
    # 假设模板文件在 backend/templates/simple_template.pptx
    TEST_TEMPLATE_PATH = os.path.join(BASE_DIR_PPPTX, "templates", "simple_template.pptx") # <--- **请检查并修改路径**
    # 或者设置为 None 不使用模板:
    # TEST_TEMPLATE_PATH = None

    print("Running PPTX generation test with local image insertion...")
    try:
        # 调用函数生成 PPTX 文件，传入数据和模板路径
        # 注意：如果TEST_TEMPLATE_PATH不存在，create_presentation_file会创建空白演示文稿
        generated_filepath = create_presentation_file(test_id, test_slides_data_with_local_image, template_path=TEST_TEMPLATE_PATH)
        print(f"Test PPTX generated at: {generated_filepath}")
        # 您可以在文件浏览器中找到这个文件并打开检查，看是否插入了图片。
    except Exception as e:
        print(f"PPTX test failed: {e}")
    finally:
        # 清理测试图片文件
        if test_image_present and os.path.exists(test_image_path):
            try:
                os.remove(test_image_path)
                print(f"Removed test image file: {test_image_path}")
            except OSError as e:
                print(f"Error removing test image file {test_image_path}: {e}")