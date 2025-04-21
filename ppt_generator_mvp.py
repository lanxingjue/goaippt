# -*- coding: utf-8 -*-

# 导入 python-pptx 库
from pptx import Presentation
from pptx.util import Inches, Pt # 用于设置尺寸和字体大小，MVP可以先不关注这个，用默认布局

# --- 模拟AI生成PPT数据 ---
# 在实际开发中，这一部分会替换为调用大型语言模型的API
def simulate_ai_generate_ppt_data(topic):
    """
    模拟AI根据主题生成PPT所需的数据结构。

    Args:
        topic (str): 用户输入的PPT主题。

    Returns:
        list: 一个包含字典的列表，每个字典代表一页幻灯片的数据。
              例如: [{'title': '...', 'content': ['...', '...'], 'notes': '...'}, ...]
              第一个字典通常用于标题页（尽管这里简化处理）。
    """
    print(f"模拟AI正在为主题 '{topic}' 生成PPT数据...")

    # 这是一个硬编码的模拟响应，根据不同的主题可以返回不同的数据
    # 实际AI应该能根据topic动态生成这些内容
    if "猫" in topic or "cat" in topic.lower():
        data = [
            {
                'title': f'{topic}：一个可爱的伙伴', # 标题页标题
                'content': [], # 标题页内容（这里简化，不填充内容placeholder）
                'notes': f'欢迎大家来听关于{topic}的分享。{topic}是人类的好朋友。' # 标题页备注
            },
            {
                'title': '猫的起源与驯化', # 幻灯片2标题
                'content': [ # 幻灯片2内容要点
                    '家猫的祖先是野猫',
                    '约在9500年前开始被人类驯化',
                    '最初可能因为捕鼠而被接纳',
                ],
                'notes': '这一页我们来简单回顾一下猫咪是如何从野外走进我们的家庭的。它们的驯化历史比我们想象的要长。' # 幻灯片2备注
            },
            {
                'title': '猫的品种多样性', # 幻灯片3标题
                'content': [ # 幻灯片3内容要点
                    '世界各地有众多猫的品种',
                    '根据体型、毛色、性格等区分',
                    '常见品种：英短、美短、布偶猫、暹罗猫等',
                    '每个品种都有其独特魅力',
                ],
                'notes': '猫咪的世界非常丰富多彩，有各种各样的品种，每一种都有其独特的外观和个性。比如布偶猫就很温柔粘人。' # 幻灯片3备注
            },
            {
                'title': '如何照顾你的猫咪', # 幻灯片4标题
                'content': [ # 幻灯片4内容要点
                    '提供均衡的食物和清水',
                    '定期梳毛和清洁',
                    '提供玩耍和运动的空间',
                    '定期体检和疫苗接种',
                    '给予足够的爱和陪伴',
                ],
                'notes': '养猫需要负责任。提供合适的饮食、干净的环境以及充分的互动，都是让猫咪健康快乐的关键。' # 幻灯片4备注
            },
            {
                'title': '总结与问答', # 幻灯片5标题
                'content': [ # 幻灯片5内容要点
                    '猫咪是可爱的家庭成员',
                    '了解它们的习性，科学喂养',
                    '享受与猫咪共度的时光',
                    '欢迎提问',
                ],
                'notes': '感谢大家的聆听！希望今天的分享能帮助大家更好地了解和照顾猫咪。现在是提问环节。' # 幻灯片5备注
            }
        ]
    elif "狗" in topic or "dog" in topic.lower():
         data = [
            {
                'title': f'{topic}：人类最忠诚的朋友', # 标题页标题
                'content': [],
                'notes': f'欢迎大家来听关于{topic}的分享。{topic}因其忠诚和友善，被誉为人类最好的朋友。' # 标题页备注
            },
            {
                'title': '狗的起源与驯化', # 幻灯片2标题
                'content': [
                    '起源于灰狼',
                    '驯化历史至少有1.5万年',
                    '与人类共同进化',
                    '早期用于狩猎、护卫等',
                ],
                'notes': '狗是人类最早驯化的动物之一。它们与人类共同经历了漫长的历史，关系非常紧密。' # 幻灯片2备注
            },
            {
                'title': '常见的狗狗品种', # 幻灯片3标题
                'content': [
                    '工作犬、伴侣犬、运动犬等分类',
                    '拉布拉多、金毛、贵宾、德牧等',
                    '选择适合自己生活方式的品种',
                ],
                'notes': '狗狗的品种非常多，每种都有独特的特点。选择狗狗时要考虑清楚自己是否有能力满足它的运动和社交需求。' # 幻灯片3备注
            },
            {
                'title': '狗狗的日常护理', # 幻灯片4标题
                'content': [
                    '均衡饮食和充足饮水',
                    '定期的运动和训练',
                    '毛发和皮肤护理',
                    '疫苗接种和驱虫',
                    '提供稳定的生活环境',
                ],
                'notes': '科学喂养和悉心照料是狗狗健康长寿的基础。别忘了花时间陪伴它们，狗狗非常需要主人的关注。' # 幻灯片4备注
            },
            {
                'title': '狗狗与人类的关系', # 幻灯片5标题
                'content': [
                    '家庭成员的角色',
                    '提供情感支持和陪伴',
                    '导盲犬、搜救犬等工作犬',
                    '增进人类健康和幸福感',
                ],
                'notes': '狗狗不仅仅是宠物，它们是我们的家人。它们在很多方面帮助和丰富着我们的生活。' # 幻灯片5备注
            }
        ]
    else: # 默认主题或其他主题
         data = [
            {
                'title': f'关于主题：“{topic}”', # 标题页标题
                'content': [],
                'notes': f'今天，我们将一起探讨关于“{topic}”这个主题。' # 标题页备注
            },
            {
                'title': '引言与背景', # 幻灯片2标题
                'content': [
                    f'介绍{topic}的基本概念',
                    f'{topic}的重要性或相关性',
                    '简述本次分享将涵盖的内容',
                ],
                'notes': f'首先，我们来了解一下{topic}的背景和它为什么值得我们关注。今天的分享会涉及以下几个方面...' # 幻灯片2备注
            },
            {
                'title': '核心观点/内容点1', # 幻灯片3标题
                'content': [
                    '第一个主要观点或事实',
                    '支撑此观点的数据或例子',
                    '进一步的解释或细节',
                ],
                'notes': '现在我们深入探讨第一个核心点。这个观点非常重要，因为它直接影响到... 这里有一些具体的数据可以佐证。' # 幻灯片3备注
            },
            {
                'title': '核心观点/内容点2', # 幻灯片4标题
                'content': [
                    '第二个主要观点或方面',
                    '与前一个观点的联系或区别',
                    '相关的案例或应用',
                ],
                'notes': '接下来看第二个要点。这与我们刚才讨论的有所不同，或者是在其基础上的延伸... 来看一个相关的实际案例。' # 幻灯片4备注
            },
            {
                'title': '结论与展望', # 幻灯片5标题
                'content': [
                    '总结主要观点',
                    '提出结论或建议',
                    '对未来的展望或思考',
                ],
                'notes': '最后，我们来总结一下今天的讨论，并对{topic}的未来做一些展望。谢谢大家的聆听，期待交流。' # 幻灯片5备注
            }
        ]

    print("模拟AI数据生成完毕。")
    return data

# --- 使用生成的数据创建PPT ---
def create_presentation_from_data(ppt_data, output_filename="generated_presentation.pptx"):
    """
    根据模拟AI生成的数据创建PPT文件。

    Args:
        ppt_data (list): 由 simulate_ai_generate_ppt_data 返回的数据列表。
        output_filename (str): 保存PPT的文件名 (默认为 "generated_presentation.pptx")。
    """
    print(f"正在创建PPT文件：{output_filename}")

    # 创建一个新的演示文稿对象
    prs = Presentation()

    # 获取常用的幻灯片布局
    # 这些索引值可能因不同的默认模板而异，但0通常是标题页，1通常是标题和内容页
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # 遍历生成的数据，创建幻灯片
    for i, slide_data in enumerate(ppt_data):
        slide_title = slide_data.get('title', '无标题')
        slide_content = slide_data.get('content', [])
        slide_notes = slide_data.get('notes', '')

        if i == 0: # 第一页通常用作标题页
            slide = prs.slides.add_slide(title_slide_layout)
            # 找到标题和副标题的placeholder（这里简化，只填充标题）
            title_shape = slide.shapes.title
            # 如果有副标题，可以尝试填充，但MVP简化，不处理
            # body_shape = slide.placeholders[1] # 通常副标题是placeholder[1]
            title_shape.text = slide_title

            # 在标题页添加备注
            if slide_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_notes

        else: # 后续页面使用标题+内容布局
            slide = prs.slides.add_slide(bullet_slide_layout)
            # 找到标题和内容的placeholder
            title_shape = slide.shapes.title
            body_shape = slide.shapes.placeholders[1] # 内容通常是placeholder[1]

            # 填充标题
            title_shape.text = slide_title

            # 填充内容要点
            tf = body_shape.text_frame
            tf.clear() # 清空默认的文本，准备添加新的段落

            for point in slide_content:
                p = tf.add_paragraph()
                p.text = point
                # 可以设置要点层级，但MVP简化，所有要点都在第一级
                # p.level = 0 # 或 1, 2, ...

            # 添加备注
            if slide_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_notes

        print(f"已创建幻灯片 {i+1}: '{slide_title}'")


    # 保存演示文稿
    prs.save(output_filename)
    print(f"PPT文件 '{output_filename}' 创建成功！")
    print(f"请在文件创建位置查找该文件。")


# --- 主程序入口 ---
if __name__ == "__main__":
    print("欢迎使用AI快速生成PPT工具（MVP版本）!")
    print("请注意：这是一个DEMO，AI生成部分是模拟的，模板和内容是固定的简单示例。")

    # 获取用户输入的主题
    user_topic = input("请输入您想要制作PPT的主题（例如：猫，狗，人工智能）：")

    # 模拟AI生成PPT数据
    ppt_data = simulate_ai_generate_ppt_data(user_topic)

    # 如果AI模拟返回了数据，则创建PPT文件
    if ppt_data:
        # 定义输出文件名
        output_file = "generated_presentation.pptx"
        # 创建PPT文件
        create_presentation_from_data(ppt_data, output_file)
    else:
        print("未能生成PPT数据，请尝试其他主题或稍后再试。")

    print("\nMVP阶段任务完成！")
    print("您现在可以在当前运行脚本的目录下找到生成的 'generated_presentation.pptx' 文件。")
    print("在后续阶段，我们将逐步添加更多功能，如更智能的内容生成、模板匹配、在线编辑等。")